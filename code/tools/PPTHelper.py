import pptx
import config as c
import tools as t
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
import os

rules_path = os.path.join(c.res_documents, 'ppt_rules.txt')
ppt_bg_path = os.path.join(c.res_pictures, 'ppt_bg.png')
laoluo_bg_path = os.path.join(c.res_pictures, 'laoluo.jpg')
last_bg_path = os.path.join(c.res_pictures, 'last.png')
story_bg_path = os.path.join(c.res_pictures, 'story.png')
ppt_file_name = os.path.join(c.outputs_documents_path, 'result.pptx')


# 厘米转英寸
def cm_to_in(cm):
    return Inches(cm / 2.54)


# 判断课件是否存在，不存在的新建一个空白
def ppt_existed(ppt_name):
    if not os.path.exists(ppt_name):
        prs = Presentation()
        prs.slide_height = cm_to_in(14.35)
        prs.slide_width = cm_to_in(25.5)
        prs.save(ppt_name)


# 模板1：只有一张图片
def model_1(prs, pic_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(pic_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))


# 模板2：只有一个标题
def model_2(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    title_box = slide.shapes.add_textbox(cm_to_in(3.89), cm_to_in(5.35), cm_to_in(17.61), cm_to_in(3.59))
    paragraph = title_box.text_frame.add_paragraph()
    paragraph.text = title
    paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph.font.size = Pt(60)
    paragraph.font.name = '微软雅黑'
    paragraph.font.color.rgb = RGBColor(255, 255, 255)


#  模板3：有字，有图片
def model_3(prs, title, pic_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    img = slide.shapes.add_picture(pic_path, cm_to_in(0), cm_to_in(0), height=cm_to_in(11.72))
    img.left = int(prs.slide_width / 2 + (prs.slide_width / 2 - img.width) / 2)
    img.top = int((prs.slide_height - img.height) / 2)
    title_box = slide.shapes.add_textbox(cm_to_in(2), cm_to_in(5.35), int(prs.slide_width / 3), cm_to_in(3.59))
    paragraph = title_box.text_frame.add_paragraph()
    paragraph.text = title
    paragraph.font.size = Pt(44)
    paragraph.font.name = '微软雅黑'
    paragraph.font.color.rgb = RGBColor(255, 255, 255)


# 模板4：两行文字，一大一小
def model_4(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    # 一级标题
    title_box_1 = slide.shapes.add_textbox(cm_to_in(1.27), cm_to_in(2.04), cm_to_in(22.86), cm_to_in(3.18))
    paragraph_1 = title_box_1.text_frame.add_paragraph()
    paragraph_1.text = title
    paragraph_1.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_1.font.size = Pt(44)
    paragraph_1.font.name = '微软雅黑'
    paragraph_1.font.color.rgb = RGBColor(255, 255, 255)
    # 二级标题
    title_box_2 = slide.shapes.add_textbox(cm_to_in(7.46), cm_to_in(6.4), cm_to_in(10.47), cm_to_in(2.39))
    paragraph_2 = title_box_2.text_frame.add_paragraph()
    paragraph_2.text = title
    paragraph_2.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_2.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_2.font.size = Pt(32)
    paragraph_2.font.name = '微软雅黑'
    paragraph_2.font.color.rgb = RGBColor(255, 255, 255)


# 模板5：一行文字，多个小标题
def model_5(prs, title, *content):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(ppt_bg_path, cm_to_in(0), cm_to_in(0), cm_to_in(25.4), cm_to_in(14.288))
    title_box_1 = slide.shapes.add_textbox(cm_to_in(1.27), cm_to_in(2.04), cm_to_in(22.86), cm_to_in(3.18))
    paragraph_1 = title_box_1.text_frame.add_paragraph()
    paragraph_1.text = title
    paragraph_1.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    paragraph_1.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    paragraph_1.font.size = Pt(44)
    paragraph_1.font.name = '微软雅黑'
    paragraph_1.font.color.rgb = RGBColor(255, 255, 255)
    # 动态构建小标题
    module_width = (prs.slide_width - cm_to_in(1.27) * 2) / len(content)
    for i in range(0, 3):
        title_box = slide.shapes.add_textbox(cm_to_in(1.27) + i * module_width, cm_to_in(6.4), module_width,
                                             cm_to_in(2.39))
        paragraph = title_box.text_frame.add_paragraph()
        paragraph.text = content[i]
        paragraph.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
        paragraph.font.size = Pt(32)
        paragraph.font.name = '微软雅黑'
        paragraph.font.color.rgb = RGBColor(255, 255, 255)


# 读取配置文件调用模板的方法
def read_rules(prs, filename):
    if os.path.exists(filename):
        with open(filename, 'r+', encoding='utf-8') as f:
            for rule in f:
                word_list = rule.replace('\n', '').split(',')
                if 'png' in rule or 'jpg' in rule:
                    if len(word_list) == 1:
                        model_1(prs, os.path.join(c.res_pictures, word_list[0]))
                    else:
                        model_3(prs, word_list[0], os.path.join(c.res_pictures, word_list[1]))
                else:
                    if len(word_list) == 1:
                        model_2(prs, word_list[0])
                    elif len(word_list) == 2:
                        model_4(prs, word_list[0], word_list[1])
                    elif len(word_list) == 4:
                        model_5(prs, word_list[0], word_list[1], word_list[2], word_list[3])


if __name__ == '__main__':
    t.is_dir_existed(c.outputs_documents_path)
    ppt_existed(ppt_file_name)
    presentation = Presentation(ppt_file_name)
    read_rules(presentation, rules_path)
    presentation.save(ppt_file_name)
